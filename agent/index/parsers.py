"""解析层 —— 把各种格式的文件变成带定位信息的文本块。

## 设计原则

**1. 优雅降级，绝不崩。**
解析库（pypdf / python-docx / openpyxl / python-pptx）都是可选依赖。
没装就把该格式标记为"不支持"并记录原因，而不是让整个索引作业挂掉。
核心索引能力只依赖标准库。

**2. 保留结构，不做"全文一坨"。**
每个块带 ``kind``（heading / paragraph / table / code）和 ``locator``
（页码 / sheet+行 / 行号）。这两样直接决定了切块质量和引用能否精确到位置。

**3. 解析失败必须可见。**
``ParsedDoc.error`` 会落库并出现在 ``index status`` 里。
静默丢弃会让用户问"我明明有这份合同，为什么搜不到"。
"""

from __future__ import annotations

import csv
import io
import json
import zipfile
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Callable

# 单文件解析上限，防止一个 500MB 的日志把内存吃光
MAX_PARSE_BYTES = 64 * 1024 * 1024
# 表格类最多抽多少行
MAX_TABLE_ROWS = 5000


@dataclass
class TextBlock:
    text: str
    locator: dict[str, Any] = field(default_factory=dict)
    kind: str = "paragraph"        # heading | paragraph | table | code | caption
    heading_path: list[str] = field(default_factory=list)


@dataclass
class ParsedDoc:
    blocks: list[TextBlock] = field(default_factory=list)
    title: str | None = None
    parser: str = ""
    error: str | None = None
    warnings: list[str] = field(default_factory=list)

    @property
    def ok(self) -> bool:
        return self.error is None and bool(self.blocks)


# ----------------------------------------------------------------------
# 可选依赖探测 —— 只在真正需要时导入，缺失则标记不可用
# ----------------------------------------------------------------------
def _try(name: str) -> Any:
    try:
        return __import__(name)
    except ImportError:
        return None


def available_parsers() -> dict[str, bool]:
    """供 CLI / API 展示"哪些格式当前可解析"。"""
    return {
        "text/code/markdown": True,          # 标准库
        "csv/tsv": True,
        "json/yaml": True,
        "pdf": (_try("pypdfium2") or _try("pypdf")) is not None,
        "docx": _try("docx") is not None,
        "xlsx": _try("openpyxl") is not None,
        "pptx": _try("pptx") is not None,
        "html": _try("html.parser") is not None,
    }


# ----------------------------------------------------------------------
# 纯文本 / 代码 / Markdown
# ----------------------------------------------------------------------
_TEXT_EXT = {
    "txt", "md", "markdown", "rst", "log", "csv", "tsv", "json", "yaml", "yml",
    "toml", "ini", "cfg", "conf", "env", "sql", "sh", "bash", "ps1", "bat",
    "py", "js", "ts", "tsx", "jsx", "go", "rs", "java", "kt", "c", "h", "cpp",
    "hpp", "cs", "rb", "php", "swift", "scala", "r", "m", "lua", "vue", "svelte",
    "html", "htm", "xml", "css", "scss", "less", "gitignore", "dockerfile", "makefile",
}

_CODE_EXT = {
    "py", "js", "ts", "tsx", "jsx", "go", "rs", "java", "kt", "c", "h", "cpp",
    "hpp", "cs", "rb", "php", "swift", "scala", "lua", "sh", "bash", "ps1", "sql",
}


def _decode(raw: bytes) -> tuple[str, list[str]]:
    """按常见编码依次尝试。中文环境下 GBK 文件很多，必须覆盖。"""
    warnings: list[str] = []
    for enc in ("utf-8", "utf-8-sig", "gb18030", "big5", "latin-1"):
        try:
            return raw.decode(enc), warnings
        except UnicodeDecodeError:
            continue
    warnings.append("未能确定编码，已按 utf-8 容错解码")
    return raw.decode("utf-8", errors="replace"), warnings


def parse_text(path: Path) -> ParsedDoc:
    raw = path.read_bytes()
    text, warnings = _decode(raw)
    ext = path.suffix.lower().lstrip(".")

    # Markdown：按标题切，保留层级路径
    if ext in ("md", "markdown"):
        return _parse_markdown(text, warnings)
    if ext in ("csv", "tsv"):
        return _parse_delimited(text, path, warnings)

    kind = "code" if ext in _CODE_EXT else "paragraph"
    lines = text.splitlines()
    # 按行号分段（每 60 行一块），让 locator 能精确到行
    blocks: list[TextBlock] = []
    for start in range(0, len(lines), 60):
        chunk = "\n".join(lines[start : start + 60]).strip()
        if chunk:
            blocks.append(TextBlock(text=chunk, kind=kind, locator={"line": start + 1}))
    return ParsedDoc(blocks=blocks, parser="text", warnings=warnings,
                     title=path.stem)


def _parse_markdown(text: str, warnings: list[str]) -> ParsedDoc:
    blocks: list[TextBlock] = []
    heading_path: list[str] = []
    buf: list[str] = []
    start_line = 1
    title: str | None = None

    def flush(end_line: int) -> None:
        body = "\n".join(buf).strip()
        if body:
            blocks.append(TextBlock(text=body, kind="paragraph",
                                    locator={"line": start_line},
                                    heading_path=list(heading_path)))
        buf.clear()

    for i, line in enumerate(text.splitlines(), start=1):
        if line.lstrip().startswith("#"):
            flush(i)
            level = len(line) - len(line.lstrip("#"))
            heading = line.lstrip("#").strip()
            heading_path = heading_path[: max(0, level - 1)] + [heading]
            if title is None and level == 1:
                title = heading
            blocks.append(TextBlock(text=heading, kind="heading",
                                    locator={"line": i}, heading_path=list(heading_path)))
            start_line = i + 1
        else:
            if not buf:
                start_line = i
            buf.append(line)
    flush(len(text.splitlines()))
    return ParsedDoc(blocks=blocks, title=title, parser="markdown", warnings=warnings)


def _parse_delimited(text: str, path: Path, warnings: list[str]) -> ParsedDoc:
    """CSV/TSV：**每一块都带表头**，否则单看数据行毫无意义。"""
    delim = "\t" if path.suffix.lower() == ".tsv" else ","
    reader = csv.reader(io.StringIO(text), delimiter=delim)
    rows = []
    for i, row in enumerate(reader):
        if i > MAX_TABLE_ROWS:
            warnings.append(f"仅索引前 {MAX_TABLE_ROWS} 行")
            break
        rows.append(row)
    if not rows:
        return ParsedDoc(parser="csv", error="空文件")

    header = rows[0]
    header_line = delim.join(header)
    blocks = [TextBlock(text=header_line, kind="heading", locator={"row": 1})]
    for start in range(1, len(rows), 40):
        group = rows[start : start + 40]
        body = "\n".join(delim.join(r) for r in group)
        blocks.append(TextBlock(
            text=f"{header_line}\n{body}",       # 表头随每块一起走
            kind="table", locator={"row": start + 1},
            heading_path=[f"表头: {header_line[:120]}"],
        ))
    return ParsedDoc(blocks=blocks, parser="csv", warnings=warnings, title=path.stem)


# ----------------------------------------------------------------------
# PDF
# ----------------------------------------------------------------------
def parse_pdf(path: Path) -> ParsedDoc:
    """PDF 解析。优先 pypdfium2，回退 pypdf。

    ## 为什么优先 pypdfium2

    实测同一批 8 个 PDF（174 页）：

    | 引擎       | 耗时    | 抽出字符 | 许可证        |
    |-----------|---------|---------|--------------|
    | pypdf     | 283.5s  | 388,919 | BSD-3        |
    | pypdfium2 | **10.3s** | **400,871** | Apache-2.0/BSD |

    **快 27 倍，且抽得更全。** 它基于 Google PDFium（Chrome 的 PDF 引擎）。

    没有选 PyMuPDF 是因为许可证：PyMuPDF 同样快，但是 **AGPL-3.0**，
    商业分发需要购买授权。pypdfium2 是宽松许可，做产品没有法律负担。
    这类选择在 demo 阶段无所谓，在产品阶段是硬约束。
    """
    if _try("pypdfium2") is not None:
        doc = _parse_pdf_pdfium(path)
        # pdfium 失败时仍给 pypdf 一次机会 —— 两个引擎对畸形 PDF 的容忍度不同
        if doc.ok or _try("pypdf") is None:
            return doc

    pypdf = _try("pypdf")
    if pypdf is None:
        return ParsedDoc(parser="pdf",
                         error="未安装 PDF 解析器（pip install pypdfium2）")
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        if getattr(reader, "is_encrypted", False):
            try:
                reader.decrypt("")
            except Exception:
                return ParsedDoc(parser="pdf", error="PDF 已加密，无法解析")

        blocks: list[TextBlock] = []
        warnings: list[str] = []
        empty_pages = 0
        for pno, page in enumerate(reader.pages, start=1):
            try:
                text = (page.extract_text() or "").strip()
            except Exception as exc:  # 单页失败不该毁掉整篇
                warnings.append(f"第 {pno} 页解析失败: {exc}")
                continue
            if not text:
                empty_pages += 1
                continue
            blocks.append(TextBlock(text=text, kind="paragraph", locator={"page": pno}))

        if empty_pages:
            warnings.append(
                f"{empty_pages} 页无可抽取文本（可能是扫描件，需 OCR 才能检索）"
            )
        if not blocks:
            return ParsedDoc(parser="pdf", warnings=warnings,
                             error="整份文档无可抽取文本，疑为扫描件")

        meta = getattr(reader, "metadata", None)
        title = (getattr(meta, "title", None) if meta else None) or path.stem
        return ParsedDoc(blocks=blocks, title=title, parser="pdf", warnings=warnings)
    except Exception as exc:
        return ParsedDoc(parser="pdf", error=f"{type(exc).__name__}: {exc}")


def _parse_pdf_pdfium(path: Path) -> ParsedDoc:
    try:
        import pypdfium2 as pdfium
    except ImportError:
        return ParsedDoc(parser="pdf", error="pypdfium2 不可用")

    doc = None
    try:
        doc = pdfium.PdfDocument(str(path))
        blocks: list[TextBlock] = []
        warnings: list[str] = []
        empty = 0

        for pno in range(len(doc)):
            try:
                page = doc[pno]
                text = page.get_textpage().get_text_range().strip()
            except Exception as exc:      # 单页失败不该毁掉整篇
                warnings.append(f"第 {pno + 1} 页解析失败: {exc}")
                continue
            if not text:
                empty += 1
                continue
            blocks.append(TextBlock(text=text, kind="paragraph", locator={"page": pno + 1}))

        if empty:
            warnings.append(f"{empty} 页无可抽取文本（可能是扫描件，需 OCR 才能检索）")
        if not blocks:
            return ParsedDoc(parser="pdfium", warnings=warnings,
                             error="整份文档无可抽取文本，疑为扫描件")
        return ParsedDoc(blocks=blocks, title=path.stem, parser="pdfium", warnings=warnings)
    except Exception as exc:
        return ParsedDoc(parser="pdfium", error=f"{type(exc).__name__}: {exc}")
    finally:
        if doc is not None:
            try:
                doc.close()
            except Exception:
                pass


# ----------------------------------------------------------------------
# Office
# ----------------------------------------------------------------------
def parse_docx(path: Path) -> ParsedDoc:
    if _try("docx") is None:
        return ParsedDoc(parser="docx", error="未安装 python-docx（pip install python-docx）")
    try:
        import docx
        doc = docx.Document(str(path))
        blocks: list[TextBlock] = []
        heading_path: list[str] = []
        title: str | None = None

        for i, para in enumerate(doc.paragraphs, start=1):
            text = para.text.strip()
            if not text:
                continue
            style = (para.style.name or "").lower()
            if style.startswith("heading"):
                try:
                    level = int(style.split()[-1])
                except (ValueError, IndexError):
                    level = 1
                heading_path = heading_path[: max(0, level - 1)] + [text]
                if title is None:
                    title = text
                blocks.append(TextBlock(text=text, kind="heading",
                                        locator={"paragraph": i}, heading_path=list(heading_path)))
            else:
                blocks.append(TextBlock(text=text, kind="paragraph",
                                        locator={"paragraph": i}, heading_path=list(heading_path)))

        for ti, table in enumerate(doc.tables, start=1):
            rows = ["\t".join(c.text.strip() for c in r.cells) for r in table.rows[:MAX_TABLE_ROWS]]
            body = "\n".join(r for r in rows if r.strip())
            if body:
                blocks.append(TextBlock(text=body, kind="table",
                                        locator={"table": ti}, heading_path=list(heading_path)))

        return ParsedDoc(blocks=blocks, title=title or path.stem, parser="docx")
    except Exception as exc:
        return ParsedDoc(parser="docx", error=f"{type(exc).__name__}: {exc}")


def parse_xlsx(path: Path) -> ParsedDoc:
    if _try("openpyxl") is None:
        return ParsedDoc(parser="xlsx", error="未安装 openpyxl（pip install openpyxl）")
    try:
        import openpyxl
        wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        blocks: list[TextBlock] = []
        warnings: list[str] = []

        for sheet in wb.worksheets:
            rows = []
            for i, row in enumerate(sheet.iter_rows(values_only=True), start=1):
                if i > MAX_TABLE_ROWS:
                    warnings.append(f"工作表 {sheet.title} 仅索引前 {MAX_TABLE_ROWS} 行")
                    break
                cells = ["" if v is None else str(v) for v in row]
                if any(c.strip() for c in cells):
                    rows.append((i, " | ".join(cells).rstrip(" |")))
            if not rows:
                continue

            header = rows[0][1]
            for start in range(0, len(rows), 30):
                group = rows[start : start + 30]
                body = "\n".join(t for _, t in group)
                blocks.append(TextBlock(
                    # 表头随每块走，否则单看数据行不知道每列是什么
                    text=(body if start == 0 else f"{header}\n{body}"),
                    kind="table",
                    locator={"sheet": sheet.title, "row": group[0][0]},
                    heading_path=[f"Sheet: {sheet.title}", f"表头: {header[:120]}"],
                ))
        wb.close()
        return ParsedDoc(blocks=blocks, title=path.stem, parser="xlsx", warnings=warnings)
    except Exception as exc:
        return ParsedDoc(parser="xlsx", error=f"{type(exc).__name__}: {exc}")


def parse_pptx(path: Path) -> ParsedDoc:
    if _try("pptx") is None:
        return ParsedDoc(parser="pptx", error="未安装 python-pptx（pip install python-pptx）")
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        blocks: list[TextBlock] = []
        for sno, slide in enumerate(prs.slides, start=1):
            parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        parts.append(t)
                if getattr(shape, "has_table", False):
                    parts.extend(
                        " | ".join(c.text.strip() for c in r.cells) for r in shape.table.rows
                    )
            body = "\n".join(parts).strip()
            if body:
                blocks.append(TextBlock(text=body, kind="paragraph",
                                        locator={"slide": sno},
                                        heading_path=[f"第 {sno} 页"]))
        return ParsedDoc(blocks=blocks, title=path.stem, parser="pptx")
    except Exception as exc:
        return ParsedDoc(parser="pptx", error=f"{type(exc).__name__}: {exc}")


# ----------------------------------------------------------------------
# 分派
# ----------------------------------------------------------------------
_DISPATCH: dict[str, Callable[[Path], ParsedDoc]] = {
    "pdf": parse_pdf,
    "docx": parse_docx, "docm": parse_docx,
    "xlsx": parse_xlsx, "xlsm": parse_xlsx,
    "pptx": parse_pptx, "pptm": parse_pptx,
}

# 明确不索引的：二进制、媒体、压缩包、体积巨大的产物
SKIP_EXT = {
    "exe", "dll", "so", "dylib", "bin", "dat", "db", "sqlite", "pyc", "class", "o", "a",
    "png", "jpg", "jpeg", "gif", "bmp", "webp", "svg", "ico", "tif", "tiff", "psd",
    "mp3", "mp4", "wav", "avi", "mov", "mkv", "flac", "webm",
    "zip", "tar", "gz", "bz2", "7z", "rar", "xz",
    "woff", "woff2", "ttf", "otf", "eot",
    "doc", "xls", "ppt",          # 旧版二进制 Office，需额外工具
}


def is_supported(path: Path) -> bool:
    ext = path.suffix.lower().lstrip(".")
    if ext in SKIP_EXT:
        return False
    return ext in _DISPATCH or ext in _TEXT_EXT or ext == "" or path.name.lower() in {
        "dockerfile", "makefile", "readme", "license"
    }


def parse(path: Path) -> ParsedDoc:
    """解析单个文件。**任何异常都转成 ParsedDoc.error，绝不向上抛。**"""
    try:
        size = path.stat().st_size
        if size > MAX_PARSE_BYTES:
            return ParsedDoc(parser="skip",
                             error=f"文件过大（{size / 1e6:.0f}MB > {MAX_PARSE_BYTES / 1e6:.0f}MB）")
        if size == 0:
            return ParsedDoc(parser="skip", error="空文件")

        ext = path.suffix.lower().lstrip(".")
        if ext in SKIP_EXT:
            return ParsedDoc(parser="skip", error=f"不索引的格式: .{ext}")

        handler = _DISPATCH.get(ext)
        if handler is not None:
            return handler(path)
        if ext in _TEXT_EXT or ext == "" or path.name.lower() in {
            "dockerfile", "makefile", "readme", "license"
        }:
            return parse_text(path)
        return ParsedDoc(parser="skip", error=f"未知格式: .{ext}")
    except PermissionError:
        return ParsedDoc(parser="skip", error="无读取权限")
    except Exception as exc:
        return ParsedDoc(parser="unknown", error=f"{type(exc).__name__}: {exc}")
